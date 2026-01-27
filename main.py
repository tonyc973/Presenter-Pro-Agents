import os
import re
import json
import requests
import pandas as pd
import seaborn as sns
import matplotlib
import matplotlib.pyplot as plt
import numpy as np
from typing import List, Literal, Any
from pydantic import BaseModel, Field
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from docx import Document
from crewai import Agent, Task, Crew, Process, LLM
from crewai.tools import BaseTool
from dotenv import load_dotenv
import agentops

# --- 1. CONFIG & INIT ---
load_dotenv()

# Optional Observability
if os.getenv("AGENTOPS_API_KEY"):
    agentops.init(os.getenv("AGENTOPS_API_KEY"))

matplotlib.use('Agg')

# LLM Configuration
llm_strict = LLM(model="gpt-4o", temperature=0.1)
llm_creative = LLM(model="gpt-4o", temperature=0.7)

# --- 2. ROBUST ASSET GENERATORS (With Fallbacks) ---

def fetch_high_res_image(query: str) -> str:
    """
    Tries Serper first. 
    If that fails, returns a high-quality deterministic placeholder.
    """
    print(f"      [Asset] 🔍 Searching Image for: '{query}'")
    os.makedirs("assets", exist_ok=True)
    filename = f"assets/img_{abs(hash(query))}.jpg"

    # STRATEGY A: Google Images via Serper
    if os.getenv("SERPER_API_KEY"):
        try:
            url = "https://google.serper.dev/images"
            payload = json.dumps({"q": f"{query} 4k wallpaper corporate aesthetic", "num": 3})
            headers = {'X-API-KEY': os.getenv("SERPER_API_KEY"), 'Content-Type': 'application/json'}
            resp = requests.post(url, headers=headers, data=payload, timeout=5)
            
            if resp.status_code == 200:
                results = resp.json()
                if 'images' in results:
                    for img in results['images']:
                        # Filter out small icons
                        if img.get('width') and img.get('width') > 1200:
                            img_data = requests.get(img['imageUrl'], headers={'User-Agent': 'Mozilla/5.0'}, timeout=5).content
                            with open(filename, 'wb') as f: f.write(img_data)
                            print(f"      [Asset] ✅ Downloaded from Google: {filename}")
                            return os.path.abspath(filename)
        except Exception as e:
            print(f"      [Warning] Serper failed ({e}). Switching to Fallback.")

    # STRATEGY B: Picsum Fallback (Guaranteed Success)
    try:
        print(f"      [Asset] ⚠️ Using Placeholder Fallback for: {query}")
        # Uses the hash of the query to ensure the same topic gets the same image every time
        seed = abs(hash(query)) % 1000
        fallback_url = f"https://picsum.photos/seed/{seed}/1920/1080"
        img_data = requests.get(fallback_url, timeout=10).content
        with open(filename, 'wb') as f: f.write(img_data)
        return os.path.abspath(filename)
    except Exception as e:
        print(f"      [Error] Fallback failed: {e}")
        return None

def generate_professional_chart(description: str) -> str:
    """Generates a chart. If LLM fails, generates a fallback chart."""
    print(f"      [Asset] 📊 Generating Chart: {description}")
    os.makedirs("assets", exist_ok=True)
    final_path = f"assets/chart_{pd.Timestamp.now().strftime('%H%M%S')}.png"
    
    # STRATEGY A: AI Generated Custom Chart
    try:
        code_prompt = f"""
        Write Python code to plot: "{description}"
        RULES:
        1. Use pandas, seaborn, matplotlib.pyplot.
        2. Create realistic dummy data.
        3. Theme: sns.set_theme(style="whitegrid").
        4. Palette: 'mako'.
        5. Background: Transparent (savefig transparent=True).
        6. Fonts: White color, size 12.
        7. Remove spines.
        8. Save to '{final_path}'.
        """
        coder = LLM(model="gpt-4o-mini", temperature=0)
        resp = coder.call([{"role": "user", "content": code_prompt}])
        
        match = re.search(r"```python\s*(.*?)\s*```", resp, re.DOTALL)
        if match:
            exec(match.group(1), {"pd": pd, "sns": sns, "plt": plt, "np": np})
            if os.path.exists(final_path):
                return os.path.abspath(final_path)
    except Exception as e:
        print(f"      [Warning] AI Chart failed ({e}). Generating fallback.")

    # STRATEGY B: Safety Chart (Hardcoded)
    try:
        plt.figure(figsize=(10, 6))
        sns.set_theme(style="whitegrid")
        df = pd.DataFrame({'Category': ['A', 'B', 'C', 'D'], 'Value': [10, 20, 15, 25]})
        sns.barplot(data=df, x='Category', y='Value', palette='viridis')
        plt.savefig(final_path, transparent=True)
        plt.close()
        return os.path.abspath(final_path)
    except:
        return None

# --- 3. THE BUILDER TOOL ---

class SlideData(BaseModel):
    title: str
    key_points: List[str]
    visual_type: Literal['Cinematic', 'Data']
    visual_description: str
    speaker_notes: str

class DeckBuilderInput(BaseModel):
    slides: List[SlideData] = Field(..., description="The FULL list of slides.")

class DeckBuilderTool(BaseTool):
    name: str = "deck_builder_tool"
    description: str = "Generates the deck."
    args_schema: type[BaseModel] = DeckBuilderInput

    def _run(self, slides: List[SlideData]) -> str:
        print(f"\n🚀 [BUILDER] Generating {len(slides)} slides...")
        
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        doc = Document()
        doc.add_heading("Speaker Notes", 0)

        for i, item in enumerate(slides):
            # Handle potential dict vs object issues
            slide = SlideData(**item) if isinstance(item, dict) else item
            
            print(f"   Building Slide {i+1}: {slide.title} ({slide.visual_type})")
            
            # 1. Fetch Asset (With Fallback)
            asset_path = None
            if slide.visual_type == 'Cinematic':
                asset_path = fetch_high_res_image(slide.visual_description)
            else:
                asset_path = generate_professional_chart(slide.visual_description)

            # 2. Create Slide
            ppt_slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank
            ppt_slide.background.fill.solid()
            ppt_slide.background.fill.fore_color.rgb = RGBColor(20, 20, 30)

            # 3. Add Visuals
            if asset_path and os.path.exists(asset_path):
                if slide.visual_type == 'Cinematic':
                    # Full Bleed Image
                    ppt_slide.shapes.add_picture(asset_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
                    
                    # Glass Overlay (Dark Box)
                    left, top, w, h = Inches(0.5), Inches(1), Inches(6), Inches(5.5)
                    shape = ppt_slide.shapes.add_shape(1, left, top, w, h)
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = RGBColor(0, 0, 0)
                    shape.line.fill.background()
                    
                    # Transparency
                    fill = shape.fill._xPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill")
                    alpha = OxmlElement("a:alpha")
                    alpha.set("val", "50000") # 50% opacity
                    fill.append(alpha)

                    # Text
                    txBox = ppt_slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(5.6), Inches(5))
                    p = txBox.text_frame.add_paragraph()
                    p.text = slide.title
                    p.font.size = Pt(36)
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(255, 255, 255)
                    
                    for pt in slide.key_points:
                        p = txBox.text_frame.add_paragraph()
                        p.text = f"• {pt}"
                        p.font.size = Pt(18)
                        p.font.color.rgb = RGBColor(230, 230, 230)
                        p.space_before = Pt(10)

                else: # Data Layout
                    # Chart Left
                    ppt_slide.shapes.add_picture(asset_path, Inches(0.5), Inches(1.5), height=Inches(5))
                    
                    # Title Top
                    txBox = ppt_slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(1))
                    p = txBox.text_frame.add_paragraph()
                    p.text = slide.title.upper()
                    p.font.size = Pt(32)
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(255, 255, 255)
                    
                    # Content Right
                    txBox2 = ppt_slide.shapes.add_textbox(Inches(7.5), Inches(2), Inches(5), Inches(5))
                    for pt in slide.key_points:
                        p = txBox2.text_frame.add_paragraph()
                        p.text = f"• {pt}"
                        p.font.size = Pt(20)
                        p.font.color.rgb = RGBColor(255, 255, 255)
                        p.space_after = Pt(20)

            doc.add_heading(f"Slide {i+1}: {slide.title}", level=2)
            doc.add_paragraph(slide.speaker_notes)
            doc.add_page_break()

        prs.save("final_presentation.pptx")
        doc.save("speaker_notes.docx")
        return "SUCCESS"

# --- 4. AGENTS ---

researcher = Agent(
    role="Researcher",
    goal="Find verified facts.",
    backstory="You are a data analyst.",
    tools=[], 
    llm=llm_strict
)

writer = Agent(
    role="Strategist",
    goal="Plan a 5-slide deck.",
    backstory="You are a presentation expert. You ensure every slide has content and a visual description.",
    llm=llm_creative
)

builder = Agent(
    role="Builder",
    goal="Build the deck using 'deck_builder_tool'.",
    backstory="You simply pass the plan to the tool.",
    tools=[DeckBuilderTool()],
    llm=llm_strict
)

# --- 5. EXECUTION ---

class PresentationPlan(BaseModel):
    slides: List[SlideData]

task_research = Task(
    description="Research {topic} trends.",
    expected_output="5 facts.",
    agent=researcher
)

task_plan = Task(
    description="""
    Create a 5-Slide Plan based on research.
    1. Title (Cinematic)
    2. Context (Cinematic)
    3. Market Data (Data)
    4. Strategy (Cinematic)
    5. Vision (Cinematic)
    
    Provide 'key_points' and 'visual_description' for EACH slide.
    """,
    expected_output="A filled PresentationPlan.",
    agent=writer,
    context=[task_research],
    output_pydantic=PresentationPlan
)

task_build = Task(
    description="Call 'deck_builder_tool' with the plan.",
    expected_output="Confirmation.",
    agent=builder,
    context=[task_plan]
)

if __name__ == "__main__":
    crew = Crew(
        agents=[researcher, writer, builder],
        tasks=[task_research, task_plan, task_build],
        process=Process.sequential,
        verbose=True
    )

    crew.kickoff(inputs={"topic": "The Future of AI 2030"})
    print("\n\n✅ DONE! Check 'final_presentation.pptx'")